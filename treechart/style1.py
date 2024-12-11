from collections import deque
from pptx import Presentation
from pptx.util import Inches
from core.data import Data
from treechart.classic_shapes import *
from treechart.horizontal.layout import *

class TreeChartStyle1:
    def __init__(self, intput_file:str):
        self.data = Data(intput_file)

        self.prs = Presentation()   
        self.prs.slide_width = Inches(SLIDE_WIDTH)
        self.prs.slide_height = Inches(SLIDE_HEIGHT)

        # Define positions for treechart elements
        self.layout = Layout()
        
        self.draw_ppt()
  
    def save_ppt(self, save_path:str):
        self.prs.save(save_path)
        print(f"PowerPoint saved to {save_path}")

    def draw_ppt(self):
        """
        Create the layout for the PowerPoint treechart slide.
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        root = self.data.parse_file()
        layout_boxes, layout_lines = self.layout.create_layout(root)

        queue = deque([root])  
        
        for _ in range(self.data.nodes_num):
            if not queue:
                break
            node = queue.popleft()
            DrawingBox(slide, node, layout_boxes[node.id], node.level)
            for i in range(len(layout_lines[node.id])):
                DrawingLine(slide, node, layout_lines[node.id][i][0], 
                            layout_lines[node.id][i][1])

            for child in node.children:
                queue.append(child)

if __name__ == "__main__":
    pass